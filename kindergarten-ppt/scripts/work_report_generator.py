#!/usr/bin/env python3
"""
年终工作总结PPT生成器
为幼儿园教师生成专业的年终工作汇报PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

class WorkReportGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.setup_slide_size()
    
    def setup_slide_size(self):
        """设置幻灯片尺寸为16:9"""
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
    
    def create_title_slide(self, title, teacher_name, class_name, year):
        """创建封面页"""
        slide_layout = self.prs.slide_layouts[0]  # 标题幻灯片布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 102, 153)
        
        # 设置副标题
        subtitle_shape = slide.placeholders[1]
        subtitle_text = f"{teacher_name}\n{class_name}\n{year}年度"
        subtitle_shape.text = subtitle_text
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        
        return slide
    
    def create_content_slide(self, title, content_list):
        """创建内容页"""
        slide_layout = self.prs.slide_layouts[1]  # 标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # 设置内容
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, content in enumerate(content_list):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = f"• {content}"
            p.font.size = Pt(20)
            p.space_after = Pt(12)
        
        return slide
    
    def create_summary_slide(self, achievements, improvements, plans):
        """创建总结页"""
        slide_layout = self.prs.slide_layouts[5]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "工作总结"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 成果展示
        achievements_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(4))
        achievements_frame = achievements_box.text_frame
        achievements_frame.text = "主要成果\n" + "\n".join([f"✓ {item}" for item in achievements])
        
        # 改进措施
        improvements_box = slide.shapes.add_textbox(Inches(4.5), Inches(2), Inches(4), Inches(4))
        improvements_frame = improvements_box.text_frame
        improvements_frame.text = "改进措施\n" + "\n".join([f"→ {item}" for item in improvements])
        
        # 未来计划
        plans_box = slide.shapes.add_textbox(Inches(8.5), Inches(2), Inches(4), Inches(4))
        plans_frame = plans_box.text_frame
        plans_frame.text = "未来计划\n" + "\n".join([f"★ {item}" for item in plans])
        
        return slide
    
    def generate_report(self, teacher_info, work_data):
        """生成完整的工作报告PPT"""
        # 封面页
        self.create_title_slide(
            f"{teacher_info['class_name']}年终工作总结",
            teacher_info['name'],
            teacher_info['class_name'],
            teacher_info['year']
        )
        
        # 目录页
        self.create_content_slide("汇报内容", [
            "工作回顾与总结",
            "教学成果展示", 
            "班级管理成效",
            "家园共育工作",
            "问题反思与改进",
            "下年度工作计划"
        ])
        
        # 工作回顾
        self.create_content_slide("工作回顾", work_data.get('work_review', []))
        
        # 教学成果
        self.create_content_slide("教学成果", work_data.get('teaching_results', []))
        
        # 班级管理
        self.create_content_slide("班级管理", work_data.get('class_management', []))
        
        # 家园共育
        self.create_content_slide("家园共育", work_data.get('home_school_cooperation', []))
        
        # 总结页
        self.create_summary_slide(
            work_data.get('achievements', []),
            work_data.get('improvements', []),
            work_data.get('future_plans', [])
        )
        
        return self.prs
    
    def save_presentation(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)
        print(f"PPT已保存为: {filename}")

def main():
    generator = WorkReportGenerator()
    
    # 示例数据
    teacher_info = {
        'name': '张老师',
        'class_name': '大班(1)班',
        'year': '2024'
    }
    
    work_data = {
        'work_review': [
            '完成了全年教学计划，开展主题活动12个',
            '组织班级特色活动15次，获得家长好评',
            '参与园内培训8次，提升专业技能',
            '配合园内各项工作，积极参与团队协作'
        ],
        'teaching_results': [
            '幼儿语言表达能力显著提升',
            '数学逻辑思维得到良好发展',
            '艺术创造力和想象力不断增强',
            '社交能力和合作意识明显改善'
        ],
        'class_management': [
            '建立了良好的班级常规和秩序',
            '创设了温馨舒适的班级环境',
            '制定了个性化的幼儿发展计划',
            '建立了完善的安全管理制度'
        ],
        'achievements': ['教学质量优秀', '家长满意度高', '幼儿发展良好'],
        'improvements': ['加强个别化指导', '丰富教学方法', '提升专业技能'],
        'future_plans': ['创新教学模式', '深化家园合作', '促进全面发展']
    }
    
    # 生成PPT
    ppt = generator.generate_report(teacher_info, work_data)
    generator.save_presentation('年终工作总结.pptx')

if __name__ == "__main__":
    main()