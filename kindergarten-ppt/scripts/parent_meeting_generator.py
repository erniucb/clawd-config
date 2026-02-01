#!/usr/bin/env python3
"""
家长会PPT生成器
为幼儿园教师生成专业的家长会演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class ParentMeetingGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.setup_slide_size()
    
    def setup_slide_size(self):
        """设置幻灯片尺寸"""
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
    
    def create_welcome_slide(self, class_name, teacher_name, date):
        """创建欢迎页"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = f"{class_name}家长会"
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 102, 153)
        
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = f"欢迎各位家长！\n\n主讲：{teacher_name}\n时间：{date}"
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        
        return slide
    
    def create_agenda_slide(self):
        """创建会议议程页"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "会议议程"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        agenda_items = [
            "班级情况介绍",
            "本学期教学工作汇报",
            "幼儿发展情况分析",
            "家园共育工作交流",
            "下学期工作计划",
            "家长提问与交流"
        ]
        
        for i, item in enumerate(agenda_items):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = f"{i+1}. {item}"
            p.font.size = Pt(22)
            p.space_after = Pt(15)
        
        return slide
    
    def create_class_info_slide(self, class_data):
        """创建班级信息页"""
        slide_layout = self.prs.slide_layouts[5]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "班级基本情况"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 班级信息
        info_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(4))
        info_frame = info_box.text_frame
        info_text = f"""班级信息：
        
👥 班级人数：{class_data.get('total_students', 0)}人
👦 男孩：{class_data.get('boys', 0)}人
👧 女孩：{class_data.get('girls', 0)}人
👩‍🏫 教师：{class_data.get('teachers', [])}
🏫 教室位置：{class_data.get('classroom', '')}"""
        
        info_frame.text = info_text
        info_frame.paragraphs[0].font.size = Pt(18)
        
        # 班级特色
        feature_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5), Inches(4))
        feature_frame = feature_box.text_frame
        feature_text = "班级特色：\n\n" + "\n".join([f"⭐ {feature}" for feature in class_data.get('features', [])])
        feature_frame.text = feature_text
        feature_frame.paragraphs[0].font.size = Pt(18)
        
        return slide
    
    def create_development_slide(self, development_data):
        """创建幼儿发展情况页"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "幼儿发展情况"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # 各领域发展情况
        domains = [
            ("🗣️ 语言发展", development_data.get('language', [])),
            ("🧮 数学认知", development_data.get('math', [])),
            ("🎨 艺术创造", development_data.get('art', [])),
            ("🤝 社会交往", development_data.get('social', [])),
            ("💪 身体发展", development_data.get('physical', []))
        ]
        
        for domain, achievements in domains:
            p = text_frame.add_paragraph() if text_frame.paragraphs else text_frame.paragraphs[0]
            p.text = domain
            p.font.size = Pt(20)
            p.font.bold = True
            p.space_after = Pt(8)
            
            for achievement in achievements[:2]:  # 只显示前2个要点
                p = text_frame.add_paragraph()
                p.text = f"  • {achievement}"
                p.font.size = Pt(16)
                p.space_after = Pt(5)
        
        return slide
    
    def create_cooperation_slide(self, cooperation_suggestions):
        """创建家园共育页"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "家园共育建议"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, suggestion in enumerate(cooperation_suggestions):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = f"💡 {suggestion}"
            p.font.size = Pt(18)
            p.space_after = Pt(12)
        
        return slide
    
    def create_qa_slide(self):
        """创建问答环节页"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "家长提问与交流"
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 20, 60)
        
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = """
        🙋‍♀️ 欢迎家长提问
        💬 共同交流育儿心得
        🤝 携手促进孩子成长
        
        感谢您的参与！
        """
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide
    
    def generate_parent_meeting_ppt(self, meeting_data):
        """生成完整的家长会PPT"""
        # 欢迎页
        self.create_welcome_slide(
            meeting_data['class_name'],
            meeting_data['teacher_name'],
            meeting_data['date']
        )
        
        # 议程页
        self.create_agenda_slide()
        
        # 班级情况
        self.create_class_info_slide(meeting_data['class_info'])
        
        # 教学工作汇报
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "本学期教学工作"
        content_shape = slide.placeholders[1]
        content_shape.text = "\n".join([f"📚 {item}" for item in meeting_data.get('teaching_work', [])])
        
        # 幼儿发展情况
        self.create_development_slide(meeting_data['development'])
        
        # 家园共育
        self.create_cooperation_slide(meeting_data.get('cooperation_suggestions', []))
        
        # 下学期计划
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "下学期工作计划"
        content_shape = slide.placeholders[1]
        content_shape.text = "\n".join([f"🎯 {item}" for item in meeting_data.get('next_plans', [])])
        
        # 问答环节
        self.create_qa_slide()
        
        return self.prs
    
    def save_presentation(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)
        print(f"家长会PPT已保存为: {filename}")

def main():
    generator = ParentMeetingGenerator()
    
    # 示例数据
    meeting_data = {
        'class_name': '大班(1)班',
        'teacher_name': '张老师、李老师',
        'date': '2024年12月15日',
        'class_info': {
            'total_students': 28,
            'boys': 15,
            'girls': 13,
            'teachers': ['张老师（主班）', '李老师（配班）'],
            'classroom': '教学楼二楼',
            'features': ['阅读特色班', '科学探索活动', '艺术创作工坊']
        },
        'teaching_work': [
            '完成五大领域教学目标',
            '开展主题活动8个',
            '组织户外活动每日2小时',
            '进行个别化教育指导'
        ],
        'development': {
            'language': ['词汇量显著增加', '表达能力提升'],
            'math': ['数概念清晰', '逻辑思维发展'],
            'art': ['创造力丰富', '动手能力强'],
            'social': ['合作意识增强', '交往能力提高'],
            'physical': ['大肌肉发展良好', '精细动作协调']
        },
        'cooperation_suggestions': [
            '坚持亲子阅读，培养阅读兴趣',
            '鼓励孩子独立完成力所能及的事情',
            '多与孩子交流，倾听他们的想法',
            '保持家园教育的一致性',
            '关注孩子的情绪变化，及时沟通'
        ],
        'next_plans': [
            '加强幼小衔接准备工作',
            '开展更多实践体验活动',
            '深化家园合作交流',
            '提升幼儿综合素质'
        ]
    }
    
    # 生成PPT
    ppt = generator.generate_parent_meeting_ppt(meeting_data)
    generator.save_presentation('家长会演示文稿.pptx')

if __name__ == "__main__":
    main()