#!/usr/bin/env python3
"""
安全教育PPT生成器
为幼儿园生成寒暑假安全教育PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class SafetyEducationGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.setup_slide_size()
    
    def setup_slide_size(self):
        """设置幻灯片尺寸"""
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
    
    def create_title_slide(self, season):
        """创建封面页"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = f"{season}假期安全教育"
        title_shape.text_frame.paragraphs[0].font.size = Pt(40)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 20, 60)
        
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = "安全第一 快乐假期\n让我们一起学习安全知识"
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        
        return slide
    
    def create_safety_topic_slide(self, topic, safety_rules, tips):
        """创建安全主题页"""
        slide_layout = self.prs.slide_layouts[5]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = topic
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.color.rgb = RGBColor(51, 102, 153)
        
        # 安全规则
        rules_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5.5), Inches(4))
        rules_frame = rules_box.text_frame
        rules_text = "安全规则：\n" + "\n".join([f"🔸 {rule}" for rule in safety_rules])
        rules_frame.text = rules_text
        rules_frame.paragraphs[0].font.size = Pt(18)
        rules_frame.paragraphs[0].font.bold = True
        
        # 小贴士
        tips_box = slide.shapes.add_textbox(Inches(6.5), Inches(2), Inches(5.5), Inches(4))
        tips_frame = tips_box.text_frame
        tips_text = "温馨提示：\n" + "\n".join([f"💡 {tip}" for tip in tips])
        tips_frame.text = tips_text
        tips_frame.paragraphs[0].font.size = Pt(18)
        tips_frame.paragraphs[0].font.bold = True
        
        return slide
    
    def create_interactive_slide(self, question, options, correct_answer):
        """创建互动问答页"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "安全知识小问答"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # 问题
        p = text_frame.paragraphs[0]
        p.text = f"问题：{question}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.space_after = Pt(20)
        
        # 选项
        for i, option in enumerate(options):
            p = text_frame.add_paragraph()
            p.text = f"{chr(65+i)}. {option}"
            p.font.size = Pt(18)
            p.space_after = Pt(12)
            if i == correct_answer:
                p.font.color.rgb = RGBColor(0, 128, 0)  # 正确答案用绿色
        
        return slide
    
    def generate_safety_ppt(self, season="寒"):
        """生成完整的安全教育PPT"""
        # 封面页
        self.create_title_slide(season)
        
        # 交通安全
        self.create_safety_topic_slide(
            "🚦 交通安全",
            [
                "过马路要走人行横道，看红绿灯",
                "不在马路上玩耍、追逐打闹",
                "乘车时要系好安全带，不把头手伸出窗外",
                "不乘坐无牌无证车辆"
            ],
            [
                "红灯停，绿灯行，黄灯等一等",
                "过马路时要左右看，确保安全再通过",
                "在车内要安静，不影响司机开车",
                "下车时要注意后方来车"
            ]
        )
        
        # 居家安全
        self.create_safety_topic_slide(
            "🏠 居家安全",
            [
                "不玩火，不碰电器插座",
                "不攀爬阳台、窗台等危险地方",
                "不随意开门给陌生人",
                "使用剪刀等尖锐物品要小心"
            ],
            [
                "发现火情立即告诉大人",
                "湿手不碰电器开关",
                "独自在家时要锁好门",
                "使用工具后要及时收好"
            ]
        )
        
        # 外出安全
        self.create_safety_topic_slide(
            "🌳 外出安全",
            [
                "不跟陌生人走，不吃陌生人给的东西",
                "外出时要告诉家长去哪里",
                "在人多的地方不要乱跑",
                "记住家长的电话号码"
            ],
            [
                "走失时要找警察叔叔帮忙",
                "不要一个人去偏僻的地方",
                "和家长走散时在原地等待",
                "学会大声呼救"
            ]
        )
        
        # 饮食安全
        self.create_safety_topic_slide(
            "🍎 饮食安全",
            [
                "不吃过期变质的食物",
                "饭前便后要洗手",
                "不暴饮暴食，少吃零食",
                "不喝生水，不吃不洁食物"
            ],
            [
                "看清食品保质期",
                "多吃蔬菜水果，营养均衡",
                "不买路边摊的食物",
                "有不舒服要及时告诉大人"
            ]
        )
        
        # 互动问答
        self.create_interactive_slide(
            "过马路时应该怎么做？",
            [
                "直接跑过去",
                "看红绿灯，走人行横道",
                "跟着别人走",
                "随便什么时候都可以过"
            ],
            1  # 正确答案是B
        )
        
        # 总结页
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "安全记心中，快乐过假期！"
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 20, 60)
        
        content_shape = slide.placeholders[1]
        content_shape.text = """
        🌟 安全是最重要的！
        🌟 遇到危险要冷静！
        🌟 及时求助很重要！
        🌟 快乐安全过假期！
        
        祝小朋友们假期愉快！🎉
        """
        content_shape.text_frame.paragraphs[0].font.size = Pt(24)
        content_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return self.prs
    
    def save_presentation(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)
        print(f"安全教育PPT已保存为: {filename}")

def main():
    generator = SafetyEducationGenerator()
    
    # 生成寒假安全教育PPT
    ppt = generator.generate_safety_ppt("寒")
    generator.save_presentation('寒假安全教育.pptx')
    
    # 也可以生成暑假版本
    # ppt = generator.generate_safety_ppt("暑")
    # generator.save_presentation('暑假安全教育.pptx')

if __name__ == "__main__":
    main()